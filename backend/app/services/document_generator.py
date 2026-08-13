"""
对话内 PPT 生成：从 LLM 回复中提取 ```doc_gen``` 块，生成真实 .pptx 文件。
"""
import json
import os
import uuid
import re


def extract_and_generate_ppt(llm_response: str) -> dict | None:
    """
    从 LLM 回复中提取 ```doc_gen``` JSON 块，生成 .pptx 文件。
    返回 {"filename": ..., "download_url": ...} 或 None（未找到块时）。
    """
    match = re.search(r'```doc_gen\s*\n(.*?)\n```', llm_response, re.DOTALL)
    if not match:
        return None

    try:
        data = json.loads(match.group(1))
    except json.JSONDecodeError:
        return None

    if data.get("type") != "ppt":
        return None

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    os.makedirs("outputs", exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x1E, 0x40, 0xAF)
    DARK = RGBColor(0x1F, 0x29, 0x37)

    for s in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

        # 标题
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1))
        tb.text_frame.text = s.get("title", "")
        p = tb.text_frame.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE

        # 分隔线
        line = slide.shapes.add_shape(1, Inches(1), Inches(1.5), Inches(4), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.fill.background()

        # 要点
        bullets = s.get("bullets", [])
        if bullets:
            cb = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(4.5))
            tf = cb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                pp.text = f"• {b}"
                pp.font.size = Pt(22)
                pp.font.color.rgb = DARK
                if i > 0:
                    pp.space_before = Pt(12)

        # 备注
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]

    filename = f"ppt_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join("outputs", filename)
    prs.save(filepath)

    return {"filename": filename, "download_url": f"/api/documents/download/{filename}"}
