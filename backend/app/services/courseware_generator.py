import os
import json
import tempfile
import asyncio
from typing import Dict, Any, Optional

from sqlalchemy.orm import Session


async def generate_courseware_file(
    db: Session,
    agent_id: int,
    skill_name: str,
    course_name: str,
    system_prompt: str,
    topic: str,
    audience: str,
    requirements: str,
    format_type: str,
    user_name: str,
    config: Dict[str, Any],
) -> str:
    """AI 生成课件的主入口：RAG 检索 → LLM 生成内容 → 写入 Word/PPT 文件。"""

    # v3.0：从 config 提取五层经验数据，注入课件生成
    five_layer = config.get("fiveLayerKnowledge") or config.get("five_layer_knowledge") or {}
    teaching_config = _build_teaching_config(five_layer)

    # Step 1: RAG 检索知识库
    rag_context = ""
    try:
        from .rag import retrieve_for_rag
        chunks = retrieve_for_rag(
            db=db,
            user_id=None,
            agent_id=agent_id,
            query=topic,
            top_k=5,
            similarity_threshold=0.5,
        )
        if chunks:
            rag_parts = []
            for c in chunks:
                source = c.get("filename", "未知来源")
                content = c.get("content", "")[:300]
                rag_parts.append(f"[来源: {source}]\n{content}")
            rag_context = "\n\n".join(rag_parts)
    except Exception:
        pass

    # Step 2: 调用 LLM 生成课件内容
    courseware_content = await _call_llm_for_courseware(
        system_prompt=system_prompt,
        topic=topic,
        audience=audience,
        requirements=requirements,
        rag_context=rag_context,
        course_name=course_name,
        format_type=format_type,
        config=config,
        teaching_config=teaching_config,
    )

    # Step 3: 根据格式生成文件
    if format_type == "word":
        return _build_word_courseware(
            courseware_content=courseware_content,
            topic=topic,
            course_name=course_name,
            skill_name=skill_name,
            user_name=user_name,
        )
    else:
        return _build_ppt_courseware(
            courseware_content=courseware_content,
            topic=topic,
            course_name=course_name,
            skill_name=skill_name,
            user_name=user_name,
        )


def _build_teaching_config(five_layer: Dict[str, Any]) -> str:
    """v3.0：从五层经验数据构建教学策略上下文，注入课件生成提示词。

    读取 L1-L4 数据，组织为「教学环节编排」上下文。
    """
    if not five_layer or not isinstance(five_layer, dict):
        return ""

    parts: list[str] = []

    # L1 知识体系层：知识点与重难点
    knowledge = five_layer.get("knowledge_layer") or {}
    topics = knowledge.get("topics") if isinstance(knowledge, dict) else None
    if topics and isinstance(topics, list) and len(topics) > 0:
        kp_lines = []
        for t in topics[:5]:
            if not isinstance(t, dict):
                continue
            name = t.get("name", "")
            key_points = t.get("key_points", [])
            difficulties = t.get("difficulties", [])
            line = f"- 知识点：{name}"
            if key_points and isinstance(key_points, list):
                line += f"；重点：{('、'.join(key_points[:3]))}"
            if difficulties and isinstance(difficulties, list):
                diff_texts = [d.get("point", "") if isinstance(d, dict) else str(d) for d in difficulties[:2]]
                line += f"；难点：{('、'.join(diff_texts))}"
            kp_lines.append(line)
        if kp_lines:
            parts.append("## 已沉淀知识点（L1 知识体系层）\n" + "\n".join(kp_lines))

    # L2 学生诊断层：学生痛点
    diagnosis = five_layer.get("diagnosis_layer") or {}
    pain_points = diagnosis.get("pain_points") if isinstance(diagnosis, dict) else None
    if pain_points and isinstance(pain_points, list) and len(pain_points) > 0:
        pp_lines = []
        for pp in pain_points[:3]:
            if not isinstance(pp, dict):
                continue
            pp_lines.append(
                f"- {pp.get('topic', '')}：学生表现为「{pp.get('surface_error', '')}」"
                f"，深层原因是「{pp.get('root_cause', '')}」，建议对策：{pp.get('solution', '')}"
            )
        if pp_lines:
            parts.append("## 学生常见诊断（L2 学生诊断层）\n" + "\n".join(pp_lines))

    # L3 教学策略层：教学决策步骤
    strategy = five_layer.get("strategy_layer") or {}
    strategies = strategy.get("strategies") if isinstance(strategy, dict) else None
    if strategies and isinstance(strategies, list) and len(strategies) > 0:
        st_lines = []
        for s in strategies[:3]:
            if not isinstance(s, dict):
                continue
            steps = s.get("steps", [])
            steps_text = " → ".join(steps) if isinstance(steps, list) and steps else s.get("method", "")
            st_lines.append(f"- 目标「{s.get('goal', '')}」：{steps_text}（理由：{s.get('reasoning', '')}）")
        if st_lines:
            parts.append("## 教学策略参考（L3 教学策略层）\n" + "\n".join(st_lines))

    # L4 课堂交互层：引导话术
    interaction = five_layer.get("interaction_layer") or {}
    guidance_flows = interaction.get("guidance_flows") if isinstance(interaction, dict) else None
    if guidance_flows and isinstance(guidance_flows, list) and len(guidance_flows) > 0:
        gf_lines = []
        for gf in guidance_flows[:2]:
            if not isinstance(gf, dict):
                continue
            steps = gf.get("steps", [])
            steps_text = " → ".join(steps) if isinstance(steps, list) and steps else ""
            gf_lines.append(f"- 触发「{gf.get('trigger', '')}」：{steps_text}")
        if gf_lines:
            parts.append("## 课堂引导流程（L4 课堂交互层）\n" + "\n".join(gf_lines))

    if not parts:
        return ""

    return (
        "## 教学经验沉淀（来自五层经验模型，请据此编排教学环节：导入 → 讲解 → 互动 → 总结）\n\n"
        + "\n\n".join(parts)
    )


async def _call_llm_for_courseware(
    system_prompt: str,
    topic: str,
    audience: str,
    requirements: str,
    rag_context: str,
    course_name: str,
    format_type: str,
    config: Dict[str, Any],
    teaching_config: str = "",
) -> Dict[str, Any]:
    """调用 LLM 生成课件结构化内容。"""
    import os as _os
    from openai import AsyncOpenAI

    api_key = _os.getenv("SILICONFLOW_API_KEY", "")
    base_url = _os.getenv("SILICONFLOW_BASE_URL", "https://api.siliconflow.cn/v1")
    model = _os.getenv("CHAT_MODEL", "deepseek-ai/DeepSeek-V3.2")

    if not api_key:
        return _get_fallback_content(topic, format_type, course_name)

    # 构建用户消息
    output_format_desc = ""
    if format_type == "word":
        output_format_desc = """请生成一份完整的Word教案，包含：
1. 课程标题
2. 教学目标（3-5条）
3. 教学重点与难点
4. 教学过程（按环节分：导入、新课讲解、互动练习、总结）
5. 课堂活动设计
6. 板书设计
7. 课后作业
8. 教学反思"""
    else:
        output_format_desc = """请生成一份PPT课件大纲，每个幻灯片页包含：
1. 幻灯片标题
2. 要点内容（3-5个bullet point）
3. 演讲备注（speaker notes）

幻灯片页数建议 8-15 页。"""

    doc_type = '教案' if format_type == 'word' else 'PPT课件'
    # v3.0：注入教学经验沉淀上下文
    teaching_block = (
        f"\n\n{teaching_config}\n" if teaching_config else ""
    )
    user_message = f"""请为「{course_name}」的「{topic}」生成一份{doc_type}。

面向用户：{audience or '学生'}
额外要求：{requirements or '无'}

{rag_context if rag_context else '（当前无知识库内容，请根据通用教学经验生成）'}{teaching_block}

重要：教学过程必须按「导入 → 新课讲解 → 互动练习 → 总结」四个环节编排，不要平铺要点。

{output_format_desc}

请以如下JSON格式输出：
{{
  "title": "课件标题",
  "sections": [
    {{
      "heading": "章节/幻灯片标题",
      "content": "内容或要点",
      "notes": "演讲备注或教学提示"
    }}
  ]
}}

直接输出JSON，不要Markdown包裹。"""

    # 如果 Skill 有 systemPrompt，作为系统消息的一部分
    full_system = system_prompt if system_prompt else "你是一位经验丰富的教学助手，擅长生成高质量的教学材料。"
    full_system += "\n\n请根据用户要求生成教学内容，输出必须是合法的 JSON 格式。"

    try:
        client = AsyncOpenAI(api_key=api_key, base_url=base_url)
        response = await client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": full_system},
                {"role": "user", "content": user_message},
            ],
            temperature=0.7,
            max_tokens=4096,
        )
        raw = response.choices[0].message.content.strip()

        # 尝试提取 JSON
        json_start = raw.find("{")
        json_end = raw.rfind("}") + 1
        if json_start >= 0 and json_end > json_start:
            json_str = raw[json_start:json_end]
            return json.loads(json_str)
        else:
            return _get_fallback_content(topic, format_type, course_name)

    except Exception as e:
        return _get_fallback_content(topic, format_type, course_name)


def _get_fallback_content(topic: str, format_type: str, course_name: str) -> Dict[str, Any]:
    """LLM 调用失败时的降级内容。"""
    if format_type == "word":
        return {
            "title": f"{course_name} - {topic} 教案",
            "sections": [
                {"heading": "课程标题", "content": f"《{topic}》教案", "notes": ""},
                {"heading": "教学目标", "content": f"1. 理解{topic}的核心概念\n2. 掌握{topic}的基本方法\n3. 能够应用{topic}解决实际问题", "notes": ""},
                {"heading": "教学重点与难点", "content": f"重点：{topic}的核心知识点\n难点：概念的理解与应用", "notes": ""},
                {"heading": "教学过程 - 导入", "content": "通过情境创设或问题导入，激发学生兴趣", "notes": "建议5分钟"},
                {"heading": "教学过程 - 新课讲解", "content": f"系统讲解{topic}的概念、原理和方法", "notes": "建议20分钟"},
                {"heading": "教学过程 - 互动练习", "content": "设计针对性练习，巩固所学知识", "notes": "建议10分钟"},
                {"heading": "教学过程 - 总结", "content": "梳理本节课要点，布置课后任务", "notes": "建议5分钟"},
                {"heading": "课后作业", "content": "完成相关练习题，预习下一节内容", "notes": ""},
            ],
        }
    else:
        return {
            "title": f"{course_name} - {topic}",
            "sections": [
                {"heading": f"{topic}", "content": "课程导入", "notes": "自我介绍与课程概述"},
                {"heading": "学习目标", "content": "• 理解核心概念\n• 掌握关键方法\n• 应用于实践", "notes": ""},
                {"heading": "核心概念讲解", "content": "• 定义与内涵\n• 特点与分类\n• 与相关概念的联系", "notes": "配合图示说明"},
                {"heading": "例题分析", "content": "• 典型例题1\n• 典型例题2\n• 解题思路", "notes": "引导学生思考"},
                {"heading": "课堂练习", "content": "• 基础练习\n• 进阶练习\n• 小组讨论", "notes": "巡视指导"},
                {"heading": "总结与回顾", "content": "• 本节课要点\n• 常见误区\n• 延伸思考", "notes": ""},
                {"heading": "课后作业", "content": "• 必做题\n• 选做题\n• 拓展阅读", "notes": ""},
            ],
        }


def _build_word_courseware(
    courseware_content: Dict[str, Any],
    topic: str,
    course_name: str,
    skill_name: str,
    user_name: str,
) -> str:
    """将 LLM 生成的内容写入 Word 文档。"""
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 标题
    title = doc.add_heading(courseware_content.get("title", topic), level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 元信息
    meta = doc.add_paragraph()
    meta.add_run(f'课程：{course_name} | ').font.size = Pt(9)
    meta.add_run(f'Skill：{skill_name} | ').font.size = Pt(9)
    meta.add_run(f'教师：{user_name}').font.size = Pt(9)

    doc.add_paragraph()

    # 各章节
    for i, section in enumerate(courseware_content.get("sections", [])):
        heading = section.get("heading", f"章节 {i+1}")
        content = section.get("content", "")
        notes = section.get("notes", "")

        doc.add_heading(heading, level=2)

        # 内容部分
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            p = doc.add_paragraph()
            if line.startswith(("1.", "2.", "3.", "4.", "5.", "6.", "7.", "8.", "•", "-", "*")):
                p.style = "List Bullet" if line.startswith(("•", "-", "*")) else "List Number"
                p.add_run(line.lstrip("0123456789.•-* "))
            else:
                p.add_run(line)

        # 备注/教学提示
        if notes:
            note_p = doc.add_paragraph()
            note_run = note_p.add_run(f'💡 教学提示：{notes}')
            note_run.font.size = Pt(9)
            note_run.font.color.rgb = RGBColor(128, 128, 128)
            note_run.italic = True

    # 保存（使用纯 ASCII 文件名避免 Windows 编码问题）
    tmp_dir = tempfile.gettempdir()
    import time
    filename = f'courseware_{int(time.time() * 1000) % 100000000}.docx'
    filepath = os.path.join(tmp_dir, filename)
    doc.save(filepath)
    return filepath


def _build_ppt_courseware(
    courseware_content: Dict[str, Any],
    topic: str,
    course_name: str,
    skill_name: str,
    user_name: str,
) -> str:
    """将 LLM 生成的内容写入 PPT 演示文稿。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    slide_width = Inches(13.333)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    sections = courseware_content.get("sections", [])

    # Slide 1: 封面
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.shapes.add_shape(1, 0, 0, slide_width, slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    bg.line.fill.background()

    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    tp = title_frame.paragraphs[0]
    tp.text = courseware_content.get("title", topic)
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    tp.alignment = PP_ALIGN.CENTER

    sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(1))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = f'{course_name}\n{skill_name} | {user_name}'
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    sp.alignment = PP_ALIGN.CENTER

    # 内容幻灯片
    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 左侧色条
        accent = slide.shapes.add_shape(1, 0, 0, Inches(0.3), slide_height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x26, 0x96, 0xDE)
        accent.line.fill.background()

        # 标题
        h_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        hp = h_box.text_frame.paragraphs[0]
        hp.text = section.get("heading", "")
        hp.font.size = Pt(32)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

        # 内容
        c_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(4.5))
        c_frame = c_box.text_frame
        c_frame.word_wrap = True

        content = section.get("content", "")
        lines = [l.strip() for l in content.split("\n") if l.strip()]

        for i, line in enumerate(lines):
            if i == 0:
                cp = c_frame.paragraphs[0]
            else:
                cp = c_frame.add_paragraph()

            clean = line.lstrip("0123456789.•-* ")
            if line.startswith(("•", "-", "*")):
                cp.text = f"  • {clean}"
                cp.level = 1
            elif line.startswith(tuple("0123456789")) and ". " in line[:5]:
                cp.text = f"  {clean}"
                cp.level = 1
            else:
                cp.text = clean

            cp.font.size = Pt(22)

        # 演讲备注（如果有）
        notes = section.get("notes", "")
        if notes:
            try:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    # 结束页
    end = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = end.shapes.add_shape(1, 0, 0, slide_width, slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    bg2.line.fill.background()

    e_box = end.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    ep = e_box.text_frame.paragraphs[0]
    ep.text = "谢谢聆听"
    ep.font.size = Pt(44)
    ep.font.bold = True
    ep.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    ep.alignment = PP_ALIGN.CENTER

    ep2 = e_box.text_frame.add_paragraph()
    ep2.text = f'{course_name} - {topic}'
    ep2.font.size = Pt(20)
    ep2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    ep2.alignment = PP_ALIGN.CENTER

    # 保存（使用纯 ASCII 文件名避免 Windows 编码问题）
    tmp_dir = tempfile.gettempdir()
    import time
    filename = f'courseware_{int(time.time() * 1000) % 100000000}.pptx'
    filepath = os.path.join(tmp_dir, filename)
    prs.save(filepath)
    return filepath
