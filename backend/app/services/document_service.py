"""
v2.3 文档生成服务：LLM 生成 JSON 大纲 → python-pptx/python-docx 创建真实文件。
"""
import os
import json
import uuid
import asyncio
import logging
from typing import Optional

from openai import AsyncOpenAI

from ..config import settings

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# LLM 调用
# ---------------------------------------------------------------------------

def _get_llm_client() -> AsyncOpenAI:
    return AsyncOpenAI(
        api_key=settings.SILICONFLOW_API_KEY,
        base_url=settings.SILICONFLOW_BASE_URL,
    )


async def _call_llm(prompt: str) -> str:
    """调用 DeepSeek-V3.2，返回纯文本。"""
    client = _get_llm_client()
    resp = await client.chat.completions.create(
        model=settings.CHAT_MODEL,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=4096,
    )
    return resp.choices[0].message.content.strip()


def _extract_json(text: str) -> dict:
    """从 LLM 输出中提取 JSON（兼容 markdown 代码块包裹）。"""
    # 去掉 ```json ... ``` 包裹
    if "```" in text:
        lines = text.split("\n")
        json_lines = []
        in_block = False
        for line in lines:
            if line.strip().startswith("```"):
                in_block = not in_block
                continue
            if in_block:
                json_lines.append(line)
        text = "\n".join(json_lines) if json_lines else text
    return json.loads(text)


# ---------------------------------------------------------------------------
# PPT 生成
# ---------------------------------------------------------------------------

PPT_PROMPT_TEMPLATE = """你是课程设计师。根据以下信息生成PPT大纲。

主题：{topic}
年级：{grade}
学科：{subject}
页数：{slide_count}
风格：{style}

输出JSON格式（只输出JSON，不要其他文字）：
{{
  "title": "课件标题",
  "subtitle": "副标题",
  "slides": [
    {{
      "page": 1,
      "title": "幻灯片标题",
      "type": "cover",
      "bullets": ["要点1（15-30字）", "要点2"],
      "notes": "演讲备注"
    }}
  ]
}}

要求：
- 第1页type=cover（封面）
- 第2页type=toc（目录）
- 正文每页3-5个要点，每个要点15-30字
- 最后一页type=summary（总结）
- 共{slide_count}页
- 内容用中文
"""


async def generate_ppt(
    topic: str,
    grade: str,
    subject: str,
    slide_count: int,
    style: str,
    skill_prompt: Optional[str] = None,
) -> str:
    """
    生成 PPT 文件：LLM → JSON 大纲 → python-pptx 创建 .pptx → 返回文件路径。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    # Step 1: LLM 生成大纲 JSON
    prompt = PPT_PROMPT_TEMPLATE.format(
        topic=topic, grade=grade, subject=subject,
        slide_count=slide_count, style=style,
    )
    if skill_prompt:
        prompt = f"{skill_prompt}\n\n{prompt}"

    logger.info(f"[PPT] Calling LLM for topic: {topic}")
    raw = await _call_llm(prompt)

    try:
        outline = _extract_json(raw)
    except json.JSONDecodeError:
        logger.warning("[PPT] LLM JSON parse failed, using fallback")
        outline = _ppt_fallback(topic, subject, slide_count)

    # Step 2: python-pptx 创建 .pptx
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x1A, 0x56, 0xDB)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)

    blank_layout = prs.slide_layouts[6]  # 空白布局
    slides_data = outline.get("slides", [])

    for slide_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        stype = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes", "")

        if stype == "cover":
            # 封面页：大标题居中
            left_spacer = slide.shapes.add_shape(
                1, Inches(0), Inches(0), Inches(0.15), prs.slide_height
            )
            left_spacer.fill.solid()
            left_spacer.fill.fore_color.rgb = BLUE
            left_spacer.line.fill.background()

            txBox = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.5), Inches(10), Inches(2)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = outline.get("title", topic)
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = DARK

            if outline.get("subtitle"):
                p2 = tf.add_paragraph()
                p2.text = outline["subtitle"]
                p2.font.size = Pt(20)
                p2.font.color.rgb = GRAY

        else:
            # 内容页：标题 + 分隔线 + 要点
            # 标题
            txBox = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11), Inches(1)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = BLUE

            # 分隔线
            line = slide.shapes.add_shape(
                1, Inches(0.8), Inches(1.5), Inches(11.5), Pt(3)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = BLUE
            line.line.fill.background()

            # 要点
            if bullets:
                txBox2 = slide.shapes.add_textbox(
                    Inches(1.0), Inches(2.0), Inches(11), Inches(4.5)
                )
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        bp = tf2.paragraphs[0]
                    else:
                        bp = tf2.add_paragraph()
                    bp.text = f"• {bullet}"
                    bp.font.size = Pt(22)
                    bp.font.color.rgb = DARK
                    bp.space_after = Pt(12)

        # 备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # Step 3: 保存到 outputs/
    output_dir = settings.OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)
    file_id = uuid.uuid4().hex[:8]
    filename = f"ppt_{file_id}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)
    logger.info(f"[PPT] Saved: {filepath}")
    return filepath


def _ppt_fallback(topic: str, subject: str, slide_count: int) -> dict:
    """LLM JSON 解析失败时的降级大纲。"""
    slides = [
        {"page": 1, "title": topic, "type": "cover", "bullets": [], "notes": ""},
        {"page": 2, "title": "目录", "type": "toc", "bullets": ["一、背景介绍", "二、核心内容", "三、案例分析", "四、总结"], "notes": ""},
    ]
    for i in range(3, slide_count):
        slides.append({
            "page": i, "title": f"内容 {i-2}", "type": "content",
            "bullets": [f"要点 1", f"要点 2", f"要点 3"], "notes": "",
        })
    slides.append({"page": slide_count, "title": "总结", "type": "summary", "bullets": ["回顾重点", "课后思考"], "notes": ""})
    return {"title": topic, "subtitle": subject, "slides": slides}


# ---------------------------------------------------------------------------
# Word 教案生成
# ---------------------------------------------------------------------------

WORD_PROMPT_TEMPLATE = """你是{subject}教师。根据以下信息生成教案。

主题：{topic}
年级：{grade}
课时：{duration}分钟

输出JSON格式（只输出JSON，不要其他文字）：
{{
  "title": "教案标题",
  "sections": [
    {{
      "heading": "章节标题",
      "level": 1,
      "content": "段落内容（支持多段落用\\n分隔）"
    }}
  ]
}}

必须包含以下章节：
1. 教学目标（知识目标、能力目标、素养目标）
2. 教学重点与难点
3. 教学过程（导入→讲解→练习→总结，每个环节需详细内容）
4. 板书设计
5. 教学反思

要求：
- level=1 为一级标题，level=2 为二级标题
- content 为段落正文
- 内容用中文，详实可用
"""


async def generate_word(
    topic: str,
    subject: str,
    grade: str,
    duration: str,
    skill_prompt: Optional[str] = None,
) -> str:
    """
    生成 Word 教案：LLM → JSON → python-docx 创建 .docx → 返回文件路径。
    """
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # Step 1: LLM 生成 JSON
    prompt = WORD_PROMPT_TEMPLATE.format(
        topic=topic, subject=subject, grade=grade, duration=duration,
    )
    if skill_prompt:
        prompt = f"{skill_prompt}\n\n{prompt}"

    logger.info(f"[Word] Calling LLM for topic: {topic}")
    raw = await _call_llm(prompt)

    try:
        outline = _extract_json(raw)
    except json.JSONDecodeError:
        logger.warning("[Word] LLM JSON parse failed, using fallback")
        outline = _word_fallback(topic, subject, duration)

    # Step 2: python-docx 创建 .docx
    doc = Document()

    # 设置默认字体
    style = doc.styles['Normal']
    style.font.name = '微软雅黑'
    style.font.size = Pt(11)

    # 标题
    title_para = doc.add_heading(outline.get("title", f"{topic} - 教案"), level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # 基本信息
    info_para = doc.add_paragraph()
    info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    info_run = info_para.add_run(f"学科：{subject}　年级：{grade}　课时：{duration}分钟")
    info_run.font.size = Pt(10)
    info_run.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    doc.add_paragraph("")  # 空行

    # 各章节
    for section in outline.get("sections", []):
        heading = section.get("heading", "")
        level = section.get("level", 1)
        content = section.get("content", "")

        if level <= 1:
            doc.add_heading(heading, level=1)
        else:
            doc.add_heading(heading, level=2)

        if content:
            for para_text in content.split("\n"):
                para_text = para_text.strip()
                if para_text:
                    doc.add_paragraph(para_text)

    # Step 3: 保存
    output_dir = settings.OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)
    file_id = uuid.uuid4().hex[:8]
    filename = f"doc_{file_id}.docx"
    filepath = os.path.join(output_dir, filename)
    doc.save(filepath)
    logger.info(f"[Word] Saved: {filepath}")
    return filepath


def _word_fallback(topic: str, subject: str, duration: str) -> dict:
    """LLM JSON 解析失败时的降级教案。"""
    return {
        "title": f"{topic} - 教案",
        "sections": [
            {"heading": "一、教学目标", "level": 1, "content": "知识目标：掌握本节核心概念\n能力目标：能够应用相关知识解决问题\n素养目标：培养科学思维"},
            {"heading": "二、教学重点与难点", "level": 1, "content": "重点：核心概念的理解\n难点：知识的应用与实践"},
            {"heading": "三、教学过程", "level": 1, "content": ""},
            {"heading": "1. 导入（5分钟）", "level": 2, "content": "通过实例引入主题，激发学生兴趣。"},
            {"heading": "2. 讲解（20分钟）", "level": 2, "content": "系统讲解核心知识点，结合案例说明。"},
            {"heading": "3. 练习（15分钟）", "level": 2, "content": "学生动手练习，教师巡回指导。"},
            {"heading": "4. 总结（5分钟）", "level": 2, "content": "回顾本节要点，布置课后作业。"},
            {"heading": "四、板书设计", "level": 1, "content": "左侧：核心概念\n中间：例题\n右侧：总结"},
            {"heading": "五、教学反思", "level": 1, "content": "记录课堂效果，改进方向。"},
        ],
    }
